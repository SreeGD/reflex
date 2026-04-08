#!/usr/bin/env python3
"""Generate Pulse Hackathon Pitch Deck — 7 slides, healthcare-focused."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
ACCENT_BLUE = RGBColor(0x4D, 0xAB, 0xF7)
ACCENT_GREEN = RGBColor(0x40, 0xC0, 0x57)
ACCENT_ORANGE = RGBColor(0xFF, 0xA9, 0x4D)
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_PURPLE = RGBColor(0xBE, 0x4B, 0xDB)
ACCENT_TEAL = RGBColor(0x20, 0xC9, 0x97)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xBB)
DIM = RGBColor(0x77, 0x77, 0x88)
CARD_BG = RGBColor(0x1A, 0x22, 0x3A)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def txt(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box


def mtxt(slide, l, t, w, h, lines, sz=16, color=WHITE, spacing=8):
    """Multi-line text box."""
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.font.size = Pt(sz)
        p.font.name = "Calibri"
        p.space_after = Pt(spacing)
        # Support color tuples
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.color.rgb = line[1]
            if len(line) > 2:
                p.font.bold = line[2]
        else:
            p.text = line
            p.font.color.rgb = color
    return box


def card(slide, l, t, w, h, fill=CARD_BG, text="", sz=14, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = True
        p.font.name = "Calibri"
    return shape


def divider(slide, y, color=RGBColor(0x33, 0x3A, 0x55)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(y), Inches(11.7), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ============================================================
# SLIDE 1: TITLE
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

txt(s, 1, 1.2, 11, 1.2, "PULSE", 80, ACCENT_BLUE, True, PP_ALIGN.CENTER)
txt(s, 1, 2.8, 11, 0.8, "AI-Powered Incident Management for Healthcare", 32, WHITE, False, PP_ALIGN.CENTER)
divider(s, 3.8)
txt(s, 1, 4.2, 11, 0.6, "Observe  ->  Analyze  ->  Act", 28, ACCENT_TEAL, False, PP_ALIGN.CENTER)
txt(s, 1, 5.2, 11, 0.8, "From 30-60 minute EHR outages to 10-second auto-resolution.\nClinicians stay with patients. Engineers sleep.", 20, LIGHT_GRAY, False, PP_ALIGN.CENTER)
txt(s, 1, 6.5, 11, 0.5, "Hackathon 2026", 16, DIM, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

txt(s, 0.8, 0.4, 12, 0.7, "The Problem", 38, ACCENT_RED, True)
txt(s, 0.8, 1.0, 12, 0.4, "3 AM. PagerDuty. EHR down. ER nurses writing on paper.", 20, LIGHT_GRAY)

# Left: Today
card(s, 0.8, 1.7, 5.5, 5.2, RGBColor(0x2A, 0x15, 0x15))
txt(s, 1.0, 1.8, 5, 0.5, "TODAY", 22, ACCENT_RED, True)
mtxt(s, 1.0, 2.4, 5, 4.2, [
    "30-60+ min MTTR for major incidents (industry benchmarks)",
    "8 manual steps per incident",
    "20-25% annual IT staff turnover (HIMSS)",
    "60-80% of alerts are noise",
    "",
    ("$7,900-$8,900/min", ACCENT_RED, True),
    "data center outage cost (Ponemon Institute)",
    "",
    ("Patient safety at risk:", ACCENT_ORANGE, True),
    "   No medication reconciliation",
    "   No allergy alerts",
    "   No drug interaction checks",
], 14, LIGHT_GRAY, 4)

# Right: Impact
card(s, 7.0, 1.7, 5.5, 5.2, RGBColor(0x15, 0x15, 0x2A))
txt(s, 7.2, 1.8, 5, 0.5, "THE REAL COST", 22, ACCENT_ORANGE, True)
mtxt(s, 7.2, 2.4, 5, 4.2, [
    ("$675K/week", ACCENT_RED, True),
    "in avoidable downtime + compliance risk",
    "",
    ("100-150 paper workarounds/year", WHITE, True),
    "when IT systems go down",
    "",
    ("HIPAA audit gaps", ACCENT_ORANGE, True),
    "manual incident response = incomplete trails",
    "",
    ("Knowledge walks out the door", WHITE, True),
    "tribal knowledge lost with every departure",
    "",
    ("1-3 SEV-2 incidents per week", ACCENT_RED, True),
    "across health system IT teams",
], 14, LIGHT_GRAY, 4)


# ============================================================
# SLIDE 3: THE SOLUTION
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

txt(s, 0.8, 0.4, 12, 0.7, "The Solution: Pulse", 38, ACCENT_GREEN, True)
txt(s, 0.8, 1.0, 12, 0.4, "AI pipeline that observes, analyzes, and acts — in seconds, not hours", 20, LIGHT_GRAY)

# Three pillars
pillars = [
    ("OBSERVE", "Receive alerts from\nPrometheus, ELK, OTEL\nvia standard webhooks", ACCENT_BLUE),
    ("ANALYZE", "6-node AI pipeline\nRAG + LLM + Review Agent\nUp to 95% confidence", ACCENT_ORANGE),
    ("ACT", "Auto-execute safe fixes\nDecision Brief for risky ones\nFull HIPAA audit trail", ACCENT_GREEN),
]
for i, (title, desc, color) in enumerate(pillars):
    x = 0.8 + i * 4.1
    card(s, x, 1.7, 3.7, 1.0, color, title, 22)
    txt(s, x, 2.9, 3.7, 1.2, desc, 15, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Pipeline flow
divider(s, 4.3)
txt(s, 0.8, 4.5, 12, 0.4, "AI Pipeline Flow", 20, WHITE, True)

nodes = [
    ("Intake", ACCENT_BLUE),
    ("Noise\nCheck", ACCENT_ORANGE),
    ("RCA\n(LLM+RAG)", ACCENT_ORANGE),
    ("Review\nAgent", ACCENT_BLUE),
    ("Remediation", ACCENT_GREEN),
    ("Alert", ACCENT_PURPLE),
]
for i, (label, color) in enumerate(nodes):
    x = 0.5 + i * 2.1
    card(s, x, 5.0, 1.8, 0.9, color, label, 12)

# Key differentiator
card(s, 0.8, 6.2, 12, 0.8, CARD_BG)
txt(s, 1.0, 6.3, 11, 0.6,
    "Review Agent evaluates 7 risk factors before ANY action. "
    "Patient-critical systems always require human approval with a Decision Brief. "
    "No blind auto-remediation.", 14, ACCENT_TEAL)


# ============================================================
# SLIDE 4: HOW IT WORKS (DEMO WALKTHROUGH)
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

txt(s, 0.8, 0.4, 12, 0.7, "How It Works — Live Demo", 38, ACCENT_BLUE, True)

steps = [
    ("1", "SEV-2 Alarm\nFires", "Alertmanager\nwebhook received", ACCENT_RED),
    ("2", "AI Pipeline\nRuns", "Noise check, RAG,\nRCA, risk review", ACCENT_ORANGE),
    ("3", "Decision\nMade", "Up to 95% confidence\nmedium blast radius", ACCENT_BLUE),
    ("4", "Human\nReviews", "Decision Brief:\nrisk, evidence, TTR", ACCENT_PURPLE),
    ("5", "Action\nExecuted", "One click approve\nfull audit trail", ACCENT_GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    x = 0.3 + i * 2.55
    card(s, x, 1.5, 2.3, 1.3, color, f"{title}", 15)
    txt(s, x, 2.9, 2.3, 0.8, desc, 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# What the engineer sees
divider(s, 4.0)
txt(s, 0.8, 4.2, 12, 0.4, "Decision Brief — Everything to Decide in 10 Seconds", 20, ACCENT_ORANGE, True)

brief_items = [
    ("What happened", "DB pool exhausted, connections leaked in PatientRepository"),
    ("Risk if act", "Service restarts, 30s unavailability"),
    ("Risk if wait", "Patient records inaccessible, medication checks blocked"),
    ("Evidence", "Runbook RB-101, 3 past incidents, all resolved same way"),
    ("Confidence", "95% (4 independent signals: RAG 30%, historical 30%, LLM 20%, recency 20%)"),
    ("Recommendation", "Approve restart"),
]

for i, (label, value) in enumerate(brief_items):
    row = i // 2
    col = i % 2
    x = 0.8 + col * 6.2
    y = 4.8 + row * 0.7
    txt(s, x, y, 1.5, 0.3, label + ":", 13, ACCENT_BLUE, True)
    txt(s, x + 1.6, y, 4.4, 0.3, value, 13, LIGHT_GRAY)


# ============================================================
# SLIDE 5: ARCHITECTURE & TOPOLOGY DISCOVERY
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

txt(s, 0.8, 0.4, 12, 0.7, "Architecture Discovery", 38, ACCENT_PURPLE, True)
txt(s, 0.8, 1.0, 12, 0.4, "Pulse understands your infrastructure — auto-discovered from 4 sources", 20, LIGHT_GRAY)

# Discovery sources
sources = [
    ("Config", "Service registry\n+ dependency graph", "1.0", ACCENT_GREEN),
    ("K8s Manifests", "Env vars reveal\ndependencies", "0.8", ACCENT_BLUE),
    ("Architecture Docs", "LLM extracts from\nConfluence pages", "0.7", ACCENT_ORANGE),
    ("Jira Tickets", "Incident history reveals\nfailure cascades", "0.5", ACCENT_PURPLE),
]

for i, (name, desc, weight, color) in enumerate(sources):
    x = 0.5 + i * 3.2
    card(s, x, 1.7, 2.9, 1.5, color, f"{name}\n\n{desc}", 13)
    txt(s, x, 3.3, 2.9, 0.3, f"Weight: {weight}", 11, DIM, False, PP_ALIGN.CENTER)

# Impact analysis
divider(s, 3.7)
txt(s, 0.8, 3.9, 6, 0.4, "Cascade Impact Analysis", 20, ACCENT_TEAL, True)
mtxt(s, 0.8, 4.4, 6, 2.5, [
    ("What happens if we restart pharmacy-service?", WHITE, True),
    "",
    "Upstream: medication, patient, ehr-gateway, scheduling",
    "Downstream: none (leaf service)",
    "Tier-1 services at risk: patient-service, ehr-gateway",
    "User journeys affected: patient_admission, medication_order",
    "",
    ("Propagated blast radius: LOW -> MEDIUM", ACCENT_ORANGE, True),
    "Feeds directly into Review Agent decision",
], 13, LIGHT_GRAY, 4)

# MedFlow topology (simplified)
txt(s, 7.5, 3.9, 5, 0.4, "MedFlow EHR Topology", 18, WHITE, True)
services_layout = [
    ("EHR Gateway", 9.2, 4.5, ACCENT_RED),
    ("Medication", 7.8, 5.3, ACCENT_GREEN),
    ("Scheduling", 8.8, 5.3, ACCENT_ORANGE),
    ("Patient Svc", 10.2, 5.3, ACCENT_RED),
    ("Pharmacy", 7.8, 6.1, ACCENT_GREEN),
    ("Billing", 10.2, 6.1, ACCENT_RED),
    ("Alert Svc", 9.2, 6.1, ACCENT_BLUE),
]
for label, x, y, color in services_layout:
    card(s, x, y, 1.3, 0.5, color, label, 9)


# ============================================================
# SLIDE 6: BY THE NUMBERS
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

txt(s, 0.8, 0.4, 12, 0.7, "By the Numbers", 38, ACCENT_TEAL, True)

# Before/After comparison
headers = [
    ("", 0.8, 3.0),
    ("Before", 5.5, 3.0),
    ("With Pulse", 8.5, 3.0),
    ("Impact", 11.0, 3.0),
]
for label, x, w in headers:
    txt(s, x, 1.3, w, 0.4, label, 16, ACCENT_BLUE if label else WHITE, True)

rows = [
    ("MTTR", "30-60+ min", "<10 sec", ">99%"),
    ("Manual Steps", "8 per incident", "0", "100%"),
    ("Engineer Hours/Week", "15-20 hrs", "<2 hrs", "90%"),
    ("Knowledge Retention", "0% on turnover", "100% (vector DB)", "Permanent"),
    ("Alert Noise", "100 alerts/day", "3 real incidents", "~97% (proj.)"),
    ("3 AM Pages", "Every SEV-2", "Only high-risk", "80%"),
    ("Compliance Audit", "2 weeks/quarter", "Automatic", "100%"),
    ("AI Cost per Incident", "$50-100 (engineer)", "$0.03 (LLM)", "~2000x+ ROI"),
]

for i, (metric, before, after, impact) in enumerate(rows):
    y = 1.8 + i * 0.55
    bg = CARD_BG if i % 2 == 0 else DARK_BG
    if bg == CARD_BG:
        card(s, 0.5, y - 0.05, 12.3, 0.5, bg)
    txt(s, 0.8, y, 4.5, 0.4, metric, 14, WHITE, True)
    txt(s, 5.5, y, 2.8, 0.4, before, 14, ACCENT_RED)
    txt(s, 8.5, y, 2.3, 0.4, after, 14, ACCENT_GREEN, True)
    txt(s, 11.0, y, 2, 0.4, impact, 14, ACCENT_TEAL, True)

# Healthcare-specific
divider(s, 6.3)
txt(s, 0.8, 6.5, 12, 0.6,
    "Healthcare: $7.9K-$8.9K/min outage cost (Ponemon)  |  "
    "$500K+/week avoided  |  "
    "100% HIPAA audit trail  |  "
    "Zero PHI sent to LLMs", 15, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: CHATOPS — MEET ENGINEERS WHERE THEY WORK
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

txt(s, 0.8, 0.4, 12, 0.7, "ChatOps — Meet Engineers Where They Work", 38, ACCENT_PURPLE, True)
txt(s, 0.8, 1.0, 12, 0.4, "At 3 AM, nobody opens a new dashboard. They ask questions where they already work.", 20, LIGHT_GRAY)

# Adapters
txt(s, 0.8, 1.7, 3, 0.4, "4 Adapters", 18, WHITE, True)
adapters = [
    ("Slack Bot", "Block Kit formatting\nInteractive buttons", ACCENT_ORANGE),
    ("Chat UI", "Streamlit web app\nSidebar incidents", ACCENT_PURPLE),
    ("CLI REPL", "Terminal chat\nLocal + remote", ACCENT_GREEN),
    ("REST API", "POST /chat\n12 endpoints", ACCENT_BLUE),
]
for i, (name, desc, color) in enumerate(adapters):
    y = 2.2 + i * 0.95
    card(s, 0.8, y, 1.8, 0.75, color, name, 13)
    txt(s, 2.8, y + 0.1, 2.5, 0.6, desc, 11, LIGHT_GRAY)

# 12 Tools
txt(s, 5.5, 1.7, 3, 0.4, "12 AI Tools", 18, WHITE, True)

txt(s, 5.5, 2.2, 3.5, 0.3, "Tier 1: Query", 14, ACCENT_GREEN, True)
t1 = ["search_knowledge", "query_logs", "query_metrics", "run_analysis", "get_incident", "list_incidents"]
for i, tool in enumerate(t1):
    card(s, 5.5, 2.6 + i * 0.4, 3.3, 0.32, RGBColor(0x1A, 0x3A, 0x1A), tool, 11)

txt(s, 9.3, 2.2, 3.5, 0.3, "Tier 2: Actions", 14, ACCENT_RED, True)
t2 = ["approve_action", "deny_action", "escalate", "execute_remediation"]
for i, tool in enumerate(t2):
    card(s, 9.3, 2.6 + i * 0.4, 3.3, 0.32, RGBColor(0x3A, 0x1A, 0x1A), tool, 11)

txt(s, 9.3, 4.3, 3.5, 0.3, "Topology", 14, ACCENT_BLUE, True)
t3 = ["show_topology", "analyze_impact"]
for i, tool in enumerate(t3):
    card(s, 9.3, 4.7 + i * 0.4, 3.3, 0.32, RGBColor(0x1A, 0x1A, 0x3A), tool, 11)

# Example conversations
divider(s, 5.7)
txt(s, 0.8, 5.9, 12, 0.4, "Example Conversations", 18, ACCENT_TEAL, True)

convos = [
    ('"What runbooks exist for EHR\nconnection pool issues?"', "Searches RAG knowledge\nbase, returns RB-101"),
    ('"Show me error logs\nfor patient-service"', "Queries LogsProvider,\nreturns formatted entries"),
    ('"What\'s the blast radius if\nwe restart pharmacy-service?"', "Cascade impact: 3 upstream\nTier-1 at risk, blast HIGH"),
    ('"Approve the action\nfor INC-2CD18800"', "Executes remediation,\nlogs audit trail"),
]

for i, (question, answer) in enumerate(convos):
    x = 0.5 + i * 3.2
    card(s, x, 6.4, 1.5, 0.8, ACCENT_PURPLE, question, 9)
    txt(s, x + 1.6, 6.5, 1.4, 0.7, answer, 10, LIGHT_GRAY)


# ============================================================
# SLIDE 8: EXPANDED ROADMAP
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

txt(s, 0.8, 0.4, 12, 0.7, "Roadmap — From Demo to Production", 38, ACCENT_ORANGE, True)

# Increment 1 - DONE
card(s, 0.5, 1.4, 3.0, 5.2, RGBColor(0x15, 0x2A, 0x15))
card(s, 0.5, 1.4, 3.0, 0.6, ACCENT_GREEN, "Increment 1 — DONE", 16)
mtxt(s, 0.7, 2.1, 2.6, 4.3, [
    ("Core Platform", ACCENT_GREEN, True),
    "6-node AI pipeline",
    "Review Agent (7 risk factors)",
    "Multi-signal confidence (4 signals)",
    "",
    ("ChatOps", ACCENT_PURPLE, True),
    "12 tools, Slack adapter",
    "Multi-turn conversation",
    "",
    ("Infrastructure", ACCENT_BLUE, True),
    "Webhook receiver",
    "Topology discovery (4 sources)",
    "Impact analysis",
    "",
    ("Quality", ACCENT_TEAL, True),
    "161 tests, 10 scenarios",
    "2 mock systems (ShopFast + MedFlow)",
], 11, LIGHT_GRAY, 2)

# Increment 2 - NEXT
card(s, 3.7, 1.4, 3.0, 5.2, CARD_BG)
card(s, 3.7, 1.4, 3.0, 0.6, ACCENT_ORANGE, "Increment 2 — NEXT", 16)
mtxt(s, 3.9, 2.1, 2.6, 4.3, [
    ("Real Knowledge Base", ACCENT_ORANGE, True),
    "PostgreSQL + pgvector",
    "Vector similarity search",
    "Semantic embeddings",
    "(OpenAI + local models)",
    "",
    ("Persistence", WHITE, True),
    "Incidents survive restarts",
    "Chat state persists",
    "AsyncPostgresSaver",
    "",
    ("Ingestion CLI", WHITE, True),
    "Chunk runbooks by section",
    "Embed and upsert",
    "Re-runnable pipeline",
    "",
    "Alembic migrations",
    "Graceful fallback to in-memory",
], 11, LIGHT_GRAY, 2)

# Increment 3
card(s, 6.9, 1.4, 3.0, 5.2, CARD_BG)
card(s, 6.9, 1.4, 3.0, 0.6, ACCENT_BLUE, "Increment 3", 16)
mtxt(s, 7.1, 2.1, 2.6, 4.3, [
    ("Real Infrastructure", ACCENT_BLUE, True),
    "MCP servers:",
    "  Prometheus (metrics)",
    "  Elasticsearch (logs)",
    "  Kubernetes (remediation)",
    "  Slack (notifications)",
    "  PagerDuty (escalation)",
    "",
    ("Live Discovery", WHITE, True),
    "Topology from OTEL traces",
    "K8s API integration",
    "Real-time health overlay",
    "",
    ("Slack Bot", WHITE, True),
    "Socket Mode -> Events API",
    "Production deployment",
    "Interactive approve/deny",
], 11, LIGHT_GRAY, 2)

# Increment 4
card(s, 10.1, 1.4, 3.0, 5.2, CARD_BG)
card(s, 10.1, 1.4, 3.0, 0.6, ACCENT_PURPLE, "Increment 4", 16)
mtxt(s, 10.3, 2.1, 2.6, 4.3, [
    ("Predictive Intelligence", ACCENT_PURPLE, True),
    "ML baselining",
    "Anomaly detection",
    "Prediction engine",
    "Catch incidents before",
    "they become outages",
    "",
    ("Closed Loop", WHITE, True),
    "Verify fix worked",
    "Update knowledge base",
    "Learn from every incident",
    "",
    ("Production UI", WHITE, True),
    "React frontend",
    "Role-based access",
    "Dashboard + analytics",
], 11, LIGHT_GRAY, 2)

# Timeline bar
divider(s, 6.8)
txt(s, 0.5, 6.9, 3.0, 0.3, "Mar-Apr 2026", 12, ACCENT_GREEN, True, PP_ALIGN.CENTER)
txt(s, 3.7, 6.9, 3.0, 0.3, "Apr-May 2026", 12, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
txt(s, 6.9, 6.9, 3.0, 0.3, "May-Jul 2026", 12, ACCENT_BLUE, False, PP_ALIGN.CENTER)
txt(s, 10.1, 6.9, 3.0, 0.3, "Jul-Sep 2026", 12, ACCENT_PURPLE, False, PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9: WHAT'S BUILT + ASK
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)

txt(s, 0.8, 0.3, 12, 0.6, "What We Built & What's Next", 38, ACCENT_BLUE, True)

# What's built
card(s, 0.8, 1.2, 5.5, 4.0, RGBColor(0x15, 0x2A, 0x15))
txt(s, 1.0, 1.3, 5, 0.4, "BUILT (Increment 1)", 18, ACCENT_GREEN, True)
mtxt(s, 1.0, 1.8, 5, 3.2, [
    "6-node LangGraph AI pipeline",
    "Review Agent with 7 dynamic risk factors",
    "12 ChatOps tools (query + action + topology)",
    "RAG over runbooks, Jira, Confluence",
    "Multi-signal confidence scoring (4 signals)",
    "Webhook: POST /webhook/alertmanager",
    "Multi-source topology discovery + impact analysis",
    "3 UIs: Demo, Chat, CLI + Slack adapter",
    "Switchable: ShopFast + MedFlow Healthcare",
    "161 tests, 10 scenarios, full mock system",
    ("Zero external deps: pip install and go", ACCENT_GREEN, True),
], 13, LIGHT_GRAY, 4)

# Roadmap
card(s, 7.0, 1.2, 5.5, 4.0, CARD_BG)
txt(s, 7.2, 1.3, 5, 0.4, "ROADMAP", 18, ACCENT_ORANGE, True)

road = [
    ("Increment 2", "pgvector knowledge base\nSemantic search, persistent state", ACCENT_ORANGE),
    ("Increment 3", "Real infrastructure\nPrometheus, K8s, Slack MCPs", ACCENT_BLUE),
    ("Increment 4", "Predictive intelligence\nML baselining, anomaly detection", ACCENT_PURPLE),
]
for i, (label, desc, color) in enumerate(road):
    y = 1.9 + i * 1.1
    card(s, 7.2, y, 2.0, 0.8, color, label, 13)
    txt(s, 9.4, y + 0.1, 2.8, 0.7, desc, 12, LIGHT_GRAY)

# The Ask
divider(s, 5.4)
card(s, 0.8, 5.7, 12, 1.3, RGBColor(0x1A, 0x30, 0x4A))
txt(s, 1.0, 5.8, 11.5, 0.4, "THE ASK", 20, ACCENT_TEAL, True, PP_ALIGN.CENTER)
txt(s, 1.0, 6.3, 11.5, 0.8,
    "Every minute of EHR downtime is a minute a clinician can't check a drug interaction.\n"
    "We're turning 30-60 minute incidents into 10-second non-events.\n"
    "The platform is built. The demo is live. We need your support.",
    16, WHITE, False, PP_ALIGN.CENTER)


# Save
output = "/Users/srmallip/projects/aiops/Pulse-Pitch-Deck.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
